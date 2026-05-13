#!/usr/bin/env python3
"""Hermes NAS Indexer v2 — 메타데이터 추출 통합
원본: nas_indexer.py
개선사항:
  1. meta_extractor.py 통합 (온톨로지 기반 카테고리 분류)
  2. 구조화된 메타데이터를 Qdrant payload에 저장
  3. 추출 신뢰도 추적
  4. 실행 로그 및 통계 기록
"""
import os
import sqlite3
import json
import subprocess
import hashlib
import urllib.request
import sys
import time
from pathlib import Path
from datetime import datetime

# meta_extractor 임포트
sys.path.insert(0, '/home/raspi5p/workspace/hermes')
try:
    from meta_extractor import extract_file_metadata, ONTOLOGY
except ImportError:
    print("[error] meta_extractor.py 임포트 실패", file=sys.stderr)
    ONTOLOGY = None

# =========================================================================
# 설정
# =========================================================================

QDRANT_URL = "http://localhost:6333"
OLLAMA_EMBED_URL = "http://localhost:11434/api/embeddings"
COLLECTION = "nas_ra_docs"
STATE_DB = "/home/raspi5p/workspace/n8n-stack/hermes-ra/indexer_state.db"
LOG_FILE = "/var/log/nas_indexer.log"
CHUNK_CHARS = 800
OVERLAP_CHARS = 160
BATCH_SIZE = 50

SCAN_PATHS = [
    "/mnt/nas-ra/공통자료/DHF (인허가)/",
    "/mnt/nas-ra/변경점문서/",
    "/mnt/nas-ra/회의자료/Project회의/CYAN/인허가문서/",
    "/mnt/nas-ra/회의자료/Project회의/Retrofit/",
    "/mnt/nas-ra/회의자료/Project회의/포터블 CE MDR/",
    "/mnt/nas-ra/회의자료/Project회의/주요 Project 인허가 이슈사항/",
    "/mnt/nas-ra/회의자료/Project회의/미국 방사선등록 EPRC/",
    "/mnt/nas-ra/공통자료/Standard(국제)/",
    "/mnt/nas-ra/공통자료/RA/00_기타 작업 서류/",
    "/mnt/nas-ra/공통자료/RA/02_회사 인증서/",
    "/mnt/nas-ra/공통자료/RA/03_제품별 인증서 & 성적서/",
    "/mnt/nas-ra/공통자료/RA/06_인허가 조사 및 보고자료/",
    "/mnt/nas-ra/공통자료/RA/10_자사 품질 매뉴얼, 절차서/",
    "/mnt/nas-ra/공통자료/RA/11_Audit F-up/",
    "/mnt/nas-ra/공통자료/RA/23_규제대응/",
    "/mnt/nas-ra/공통자료/RA/52_컨설팅/",
    "/mnt/nas-ra/공통자료/RA/★User Manual/",
    "/mnt/nas-ra/공통자료/RA/★Label/",
]

SUPPORTED_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".xlsx"}
SKIP_PREFIXES = ("~$", ".")

# =========================================================================
# 로깅
# =========================================================================

def log_msg(msg: str, level: str = "info"):
    """로그 메시지 기록"""
    timestamp = datetime.now().isoformat()
    log_line = f"[{timestamp}] [{level}] {msg}"
    print(log_line, flush=True)

    try:
        with open(LOG_FILE, 'a', encoding='utf-8') as f:
            f.write(log_line + '\n')
    except:
        pass

# =========================================================================
# 데이터베이스
# =========================================================================

def init_db():
    conn = sqlite3.connect(STATE_DB)
    conn.execute("""CREATE TABLE IF NOT EXISTS indexed_files (
        path TEXT PRIMARY KEY,
        mtime REAL,
        size INTEGER,
        qdrant_ids TEXT,
        indexed_at TEXT,
        metadata_json TEXT
    )""")
    conn.commit()
    return conn

def get_file_state(conn, path):
    row = conn.execute(
        "SELECT mtime, size, qdrant_ids, metadata_json FROM indexed_files WHERE path=?",
        (path,)
    ).fetchone()
    return row

def save_file_state(conn, path, mtime, size, qdrant_ids, metadata):
    conn.execute(
        "INSERT OR REPLACE INTO indexed_files (path,mtime,size,qdrant_ids,indexed_at,metadata_json) VALUES(?,?,?,?,?,?)",
        (path, mtime, size, json.dumps(qdrant_ids), datetime.now().isoformat(), json.dumps(metadata, ensure_ascii=False))
    )
    conn.commit()

# =========================================================================
# 추출 및 임베딩
# =========================================================================

def extract_text(filepath):
    """텍스트 추출 (기존 로직 유지)"""
    ext = Path(filepath).suffix.lower()
    try:
        if ext == ".pdf":
            r = subprocess.run(["pdftotext", "-q", filepath, "-"],
                               capture_output=True, text=True, timeout=60)
            return r.stdout
        elif ext == ".docx":
            import docx
            doc = docx.Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts = []
            for sn in wb.sheetnames:
                ws = wb[sn]
                parts.append(f"[{sn}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None and str(c).strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            wb.close()
            return "\n".join(parts)
        elif ext == ".doc":
            r = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt", "--outdir", "/tmp", filepath],
                capture_output=True, timeout=120
            )
            txt = f"/tmp/{Path(filepath).stem}.txt"
            if os.path.exists(txt):
                with open(txt) as f:
                    return f.read()
    except Exception as e:
        log_msg(f"Extract error {Path(filepath).name}: {e}", "warn")
    return ""

def chunk_text(text):
    """텍스트 청킹"""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_CHARS, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - OVERLAP_CHARS
    return chunks

def embed(text):
    """OLLAMA 임베딩"""
    data = json.dumps({"model": "nomic-embed-text", "prompt": text}).encode()
    req = urllib.request.Request(OLLAMA_EMBED_URL, data=data,
                                  headers={"Content-Type": "application/json"}, method="POST")
    with urllib.request.urlopen(req, timeout=60) as resp:
        return json.loads(resp.read())["embedding"]

def qdrant_upsert(points):
    """Qdrant에 저장"""
    data = json.dumps({"points": points}).encode()
    req = urllib.request.Request(
        f"{QDRANT_URL}/collections/{COLLECTION}/points?wait=true",
        data=data, headers={"Content-Type": "application/json"}, method="PUT"
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        return json.loads(resp.read())

def qdrant_delete(ids):
    """Qdrant에서 삭제"""
    if not ids:
        return
    data = json.dumps({"points_selector": {"points_list": ids}}).encode()
    req = urllib.request.Request(
        f"{QDRANT_URL}/collections/{COLLECTION}/points/delete?wait=true",
        data=data, headers={"Content-Type": "application/json"}, method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            return json.loads(resp.read())
    except:
        pass

def make_id(filepath, chunk_idx):
    """고유 ID 생성"""
    h = hashlib.sha256(f"{filepath}:{chunk_idx}".encode()).hexdigest()
    return int(h[:16], 16) % (2**63 - 1)

# =========================================================================
# 인덱싱
# =========================================================================

def index_file(conn, filepath):
    """파일 인덱싱 (메타데이터 통합)"""
    path = Path(filepath)
    try:
        stat = path.stat()
        mtime = stat.st_mtime
        size = stat.st_size
    except:
        return "error"

    # 변경 체크
    prev = get_file_state(conn, filepath)
    if prev and prev[0] == mtime and prev[1] == size:
        return "skip"

    # 텍스트 추출
    text = extract_text(filepath)
    if not text.strip():
        return "empty"

    # 청킹
    chunks = chunk_text(text)
    if not chunks:
        return "empty"

    # 메타데이터 추출 (v2 신규)
    try:
        file_metadata = extract_file_metadata(filepath)
    except:
        file_metadata = {
            "file_name": path.name,
            "file_extension": path.suffix.lower(),
            "document_category": "unknown",
            "category_confidence": 0.0
        }

    # 기존 벡터 삭제
    if prev and prev[2]:
        old_ids = json.loads(prev[2])
        qdrant_delete(old_ids)

    fname = path.name
    new_ids = []
    batch = []

    for i, chunk in enumerate(chunks):
        try:
            vector = embed(chunk)
        except Exception as e:
            log_msg(f"Embed error {fname} chunk {i}: {e}", "warn")
            continue

        pid = make_id(filepath, i)
        new_ids.append(pid)

        # Payload 구성 (메타데이터 포함)
        payload = {
            "file_path": filepath,
            "filename": fname,
            "chunk_index": i,
            "text": chunk,
            "modified_at": mtime,
            # 메타데이터 추가
            "document_category": file_metadata.get("document_category", "unknown"),
            "category_name": file_metadata.get("category_name", "Unknown"),
            "category_confidence": file_metadata.get("category_confidence", 0.0),
            "product_code": file_metadata.get("product_code"),
            "version": file_metadata.get("version"),
            "standard_references": file_metadata.get("standard_references", []),
        }

        batch.append({
            "id": pid,
            "vector": vector,
            "payload": payload
        })

        if len(batch) >= BATCH_SIZE:
            qdrant_upsert(batch)
            batch = []

    if batch:
        qdrant_upsert(batch)

    if new_ids:
        save_file_state(conn, filepath, mtime, size, new_ids, file_metadata)
        return f"ok({len(new_ids)})"

    return "empty"

def run():
    """메인 인덱싱 루프"""
    log_msg("Starting NAS indexer v2", "info")

    conn = init_db()
    stats = {"indexed": 0, "skip": 0, "empty": 0, "error": 0}
    t0 = time.time()

    for base in SCAN_PATHS:
        if not os.path.exists(base):
            log_msg(f"Skip (not mounted): {base}", "warn")
            continue

        log_msg(f"Scan: {base}", "info")

        for root, dirs, files in os.walk(base):
            dirs[:] = sorted(d for d in dirs if not d.startswith("."))
            for fname in files:
                if any(fname.startswith(p) for p in SKIP_PREFIXES):
                    continue
                ext = Path(fname).suffix.lower()
                if ext not in SUPPORTED_EXTS:
                    continue

                fpath = os.path.join(root, fname)
                try:
                    result = index_file(conn, fpath)
                    key = "indexed" if result.startswith("ok") else result
                    stats[key] = stats.get(key, 0) + 1

                    if result.startswith("ok"):
                        log_msg(f"+ {fname} [{result}]", "debug")
                except Exception as e:
                    stats["error"] += 1
                    log_msg(f"Error {fname}: {e}", "error")

    conn.close()
    elapsed = int(time.time() - t0)

    # 최종 보고
    summary = (
        f"indexed: {stats['indexed']}, skip: {stats['skip']}, "
        f"empty: {stats['empty']}, error: {stats['error']}, elapsed: {elapsed}s"
    )
    log_msg(f"Done in {elapsed}s — {summary}", "info")

    return stats

if __name__ == "__main__":
    try:
        stats = run()
        print(json.dumps(stats, ensure_ascii=False))
    except Exception as e:
        log_msg(f"Fatal error: {e}", "error")
        sys.exit(1)
