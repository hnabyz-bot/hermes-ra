#!/usr/bin/env python3
"""Hermes NAS Indexer v1 — crawl NAS docs, embed, upsert to Qdrant nas_ra_docs"""
import os, sqlite3, json, subprocess, hashlib, urllib.request, sys, time, signal
from pathlib import Path
from datetime import datetime
from contextlib import contextmanager

FILE_EXTRACT_TIMEOUT = 30  # seconds per file

@contextmanager
def time_limit(seconds):
    def handler(signum, frame):
        raise TimeoutError(f"extract timed out after {seconds}s")
    signal.signal(signal.SIGALRM, handler)
    signal.alarm(seconds)
    try:
        yield
    finally:
        signal.alarm(0)

QDRANT_URL = os.environ.get("QDRANT_URL", "http://localhost:6333")
OLLAMA_EMBED_URL = os.environ.get("OLLAMA_URL", "http://192.168.100.1:11434") + "/api/embeddings"
COLLECTION = "nas_ra_docs"
STATE_DB = os.environ.get("STATE_DB", "/opt/hermes-ra/indexer_state.db")
CHUNK_CHARS = 800    # ~500 tokens
OVERLAP_CHARS = 160  # ~100 tokens
BATCH_SIZE = 50      # points per Qdrant upsert call

SCAN_PATHS = [
    "/mnt/nas-ra/공통자료/DHF (인허가)/",
    "/mnt/nas-ra/변경점문서/",
    "/mnt/nas-ra/회의자료/Project회의/CYAN/인허가문서/",
    "/mnt/nas-ra/회의자료/Project회의/Retrofit/",
    "/mnt/nas-ra/회의자료/Project회의/포터블 CE MDR/",
    "/mnt/nas-ra/회의자료/Project회의/주요 Project 인허가 이슈사항/",
    "/mnt/nas-ra/회의자료/Project회의/미국 방사선등록 EPRC/",
    "/mnt/nas-ra/공통자료/Standard(국제)/",
    # RA 핵심 지식 베이스 (2026-05-08 추가)
    "/mnt/nas-ra/공통자료/RA/00_기타 작업 서류/",
    "/mnt/nas-ra/공통자료/RA/02_회사 인증서/",
    "/mnt/nas-ra/공통자료/RA/03_제품별 인증서 & 성적서/",
    "/mnt/nas-ra/공통자료/RA/06_인허가 조사 및 보고자료/",
    "/mnt/nas-ra/공통자료/RA/10_자사 품질 매뉴얼, 절차서/",
    "/mnt/nas-ra/공통자료/RA/11_Audit F-up/",
    "/mnt/nas-ra/공통자료/RA/23_규제대응/",
    "/mnt/nas-ra/공통자료/RA/52_컨설팅/",
    "/mnt/nas-ra/공통자료/RA/★User Manual/",
    "/mnt/nas-ra/공통자료/RA/★Label/",
]

SUPPORTED_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".xlsx"}
SKIP_PREFIXES = ("~$", ".")


def init_db():
    conn = sqlite3.connect(STATE_DB)
    conn.execute("""CREATE TABLE IF NOT EXISTS indexed_files (
        path TEXT PRIMARY KEY,
        mtime REAL,
        size INTEGER,
        qdrant_ids TEXT,
        indexed_at TEXT
    )""")
    conn.commit()
    return conn


def get_file_state(conn, path):
    row = conn.execute("SELECT mtime, size, qdrant_ids FROM indexed_files WHERE path=?", (path,)).fetchone()
    return row  # (mtime, size, qdrant_ids_json) or None


def save_file_state(conn, path, mtime, size, qdrant_ids):
    conn.execute(
        "INSERT OR REPLACE INTO indexed_files (path,mtime,size,qdrant_ids,indexed_at) VALUES(?,?,?,?,?)",
        (path, mtime, size, json.dumps(qdrant_ids), datetime.now().isoformat())
    )
    conn.commit()


def extract_text(filepath):
    ext = Path(filepath).suffix.lower()
    try:
        with time_limit(FILE_EXTRACT_TIMEOUT):
            if ext == ".pdf":
                r = subprocess.run(["pdftotext", "-q", filepath, "-"],
                                   capture_output=True, text=True, timeout=25)
                return r.stdout
            elif ext == ".docx":
                import docx
                doc = docx.Document(filepath)
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(filepath)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                return "\n".join(parts)
            elif ext == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                parts = []
                for sn in wb.sheetnames:
                    ws = wb[sn]
                    parts.append(f"[{sn}]")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None and str(c).strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                wb.close()
                return "\n".join(parts)
            elif ext in (".doc",):
                r = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "txt", "--outdir", "/tmp", filepath],
                    capture_output=True, timeout=25
                )
                txt = f"/tmp/{Path(filepath).stem}.txt"
                if os.path.exists(txt):
                    with open(txt) as f:
                        return f.read()
    except Exception as e:
        print(f"  [extract error] {Path(filepath).name}: {e}", file=sys.stderr)
    return ""


def chunk_text(text):
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_CHARS, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - OVERLAP_CHARS
    return chunks


def embed(text):
    model = os.environ.get("EMBED_MODEL", "qwen3-embedding:latest")
    data = json.dumps({"model": model, "prompt": text}).encode()
    req = urllib.request.Request(OLLAMA_EMBED_URL, data=data,
                                  headers={"Content-Type": "application/json"}, method="POST")
    with urllib.request.urlopen(req, timeout=60) as resp:
        return json.loads(resp.read())["embedding"]


def qdrant_upsert(points):
    data = json.dumps({"points": points}).encode()
    req = urllib.request.Request(
        f"{QDRANT_URL}/collections/{COLLECTION}/points?wait=true",
        data=data, headers={"Content-Type": "application/json"}, method="PUT"
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        return json.loads(resp.read())


def qdrant_delete(ids):
    if not ids:
        return
    data = json.dumps({"points": ids}).encode()
    req = urllib.request.Request(
        f"{QDRANT_URL}/collections/{COLLECTION}/points/delete",
        data=data, headers={"Content-Type": "application/json"}, method="POST"
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        return json.loads(resp.read())


def make_id(path, chunk_index):
    h = hashlib.md5(f"{path}:{chunk_index}".encode()).hexdigest()
    return int(h[:15], 16)  # 60-bit int, safe for Qdrant


def folder_category(path):
    for p in SCAN_PATHS:
        if path.startswith(p):
            return Path(p).name
    return "other"


def index_file(conn, filepath):
    try:
        stat = os.stat(filepath)
    except OSError:
        return "error"
    mtime, size = stat.st_mtime, stat.st_size

    prev = get_file_state(conn, filepath)
    if prev and abs(prev[0] - mtime) < 1 and prev[1] == size:
        return "skip"

    text = extract_text(filepath)
    if not text.strip():
        return "empty"

    chunks = chunk_text(text)
    if not chunks:
        return "empty"

    # Delete old vectors
    if prev and prev[2]:
        old_ids = json.loads(prev[2])
        qdrant_delete(old_ids)

    fname = Path(filepath).name
    category = folder_category(filepath)
    new_ids = []
    batch = []

    for i, chunk in enumerate(chunks):
        try:
            vector = embed(chunk)
        except Exception as e:
            print(f"  [embed error] chunk {i}: {e}", file=sys.stderr)
            continue
        pid = make_id(filepath, i)
        new_ids.append(pid)
        batch.append({
            "id": pid,
            "vector": vector,
            "payload": {
                "file_path": filepath,
                "filename": fname,
                "folder_category": category,
                "modified_at": mtime,
                "chunk_index": i,
                "text": chunk  # 전체 청크 저장 (RAG 품질 향상)
            }
        })
        if len(batch) >= BATCH_SIZE:
            qdrant_upsert(batch)
            batch = []

    if batch:
        qdrant_upsert(batch)

    if new_ids:
        save_file_state(conn, filepath, mtime, size, new_ids)

    return f"ok({len(new_ids)}chunks)"


def qdrant_point_count():
    try:
        req = urllib.request.Request(f"{QDRANT_URL}/collections/{COLLECTION}")
        resp = json.loads(urllib.request.urlopen(req, timeout=5).read())
        return resp.get("result", {}).get("points_count", 0)
    except Exception:
        return -1


def run(force_reindex=False):
    conn = init_db()

    # 신규 PC 안전장치: DB에 기록은 있으나 Qdrant가 비어있으면 자동 경고
    if not force_reindex:
        db_count = conn.execute("SELECT COUNT(*) FROM indexed_files").fetchone()[0]
        qdrant_count = qdrant_point_count()
        if db_count > 0 and qdrant_count == 0:
            print("[warn] ⚠️  indexer_state.db에 기록이 있으나 Qdrant가 비어있습니다.")
            print("[warn] 신규 PC 이전 상황으로 판단됩니다.")
            print("[warn] --force-reindex 로 재실행하거나 qdrant_restore.sh 로 스냅샷을 복원하세요.")
            sys.exit(1)

    if force_reindex:
        print("[reindex] indexer_state.db 초기화 후 전체 재인덱싱")
        conn.execute("DELETE FROM indexed_files")
        conn.commit()

    stats = {}
    t0 = time.time()

    for base in SCAN_PATHS:
        if not os.path.exists(base):
            print(f"[skip] not mounted: {base}")
            continue
        print(f"\n[scan] {base}")
        for root, dirs, files in os.walk(base):
            dirs[:] = sorted(d for d in dirs if not d.startswith("."))
            for fname in files:
                if any(fname.startswith(p) for p in SKIP_PREFIXES):
                    continue
                ext = Path(fname).suffix.lower()
                if ext not in SUPPORTED_EXTS:
                    continue
                fpath = os.path.join(root, fname)
                result = index_file(conn, fpath)
                key = "indexed" if result.startswith("ok") else result
                stats[key] = stats.get(key, 0) + 1
                if result.startswith("ok"):
                    print(f"  + {fname} [{result}]")

    conn.close()
    elapsed = int(time.time() - t0)
    print(f"\n=== Done in {elapsed}s ===")
    print(f"  indexed: {stats.get('indexed',0)}, skipped: {stats.get('skip',0)}, "
          f"empty: {stats.get('empty',0)}, error: {stats.get('error',0)}")


if __name__ == "__main__":
    force = "--force-reindex" in sys.argv or "--reindex" in sys.argv
    run(force_reindex=force)
