#!/usr/bin/env python3
# Hermes RA 분석 API 서버 v5.1
# 변경사항: /analyze 첨부파일 추출 + NAS RAG + 3단계 GLM 캐스케이드 + wp_comment 출력
# /qa, /health 엔드포인트 유지 (v5 동일)
from http.server import BaseHTTPRequestHandler
from socketserver import ThreadingMixIn
from http.server import HTTPServer
import json, re, urllib.request, urllib.error, os, base64, subprocess, tempfile
from pathlib import Path

# ---------------------------------------------------------------------------
# 설정
# ---------------------------------------------------------------------------

GLM_API_KEY = os.environ.get("GLM_API_KEY", "")
GLM_BASE_URL = "https://api.z.ai/api/paas/v4/chat/completions"
NAS_QDRANT_SEARCH = "http://localhost:6333/collections/nas_ra_docs/points/search"
OLLAMA_EMBED_URL = "http://localhost:11434/api/embeddings"
OLLAMA_GENERATE_URL = "http://localhost:11434/api/generate"

# /qa 엔드포인트용 (기존 hermes-ra-knowledge 컬렉션)
QA_QDRANT_SEARCH_URL = "http://localhost:6333/collections/hermes-ra-knowledge/points/search"
QA_QDRANT_COLLECTION_URL = "http://localhost:6333/collections/hermes-ra-knowledge"
QA_QDRANT_SCROLL_URL = "http://localhost:6333/collections/hermes-ra-knowledge/points/scroll"
OLLAMA_TAGS_URL = "http://localhost:11434/api/tags"

# ---------------------------------------------------------------------------
# 공통 유틸
# ---------------------------------------------------------------------------

def _post_json(url, payload, timeout=30):
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        url, data=data,
        headers={"Content-Type": "application/json"}, method="POST"
    )
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return json.loads(resp.read().decode("utf-8"))


def _get_json(url, timeout=10):
    try:
        with urllib.request.urlopen(url, timeout=timeout) as resp:
            return json.loads(resp.read().decode("utf-8"))
    except Exception:
        return None


def parse_json_response(raw):
    raw = re.sub(r"```json\s*", "", raw)
    raw = re.sub(r"```\s*", "", raw)
    raw = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL)
    raw = raw.strip()
    try:
        return json.loads(raw)
    except Exception:
        m = re.search(r"\{.*\}", raw, re.DOTALL)
        if m:
            try:
                return json.loads(m.group())
            except Exception:
                pass
    return None


def clean_body(body):
    body = re.sub(r"https?://\S+", "", body)
    body = re.sub(r"^(From|To|Cc|Subject|Date)\s*:\s*[^\n]+\n?", "", body, flags=re.MULTILINE)
    body = re.sub(r"\s+", " ", body).strip()
    return body[:5000]


# ---------------------------------------------------------------------------
# v5.1 — 첨부파일 텍스트 추출
# ---------------------------------------------------------------------------

def extract_attachment_text(attachment_files):
    """attachment_files: [{filename, content_type, data(base64)}] → 합쳐진 텍스트"""
    texts = []
    for att in attachment_files:
        fname = att.get("filename", "")
        data_b64 = att.get("data", "")
        if not data_b64:
            continue
        ext = Path(fname).suffix.lower()
        try:
            # Gmail API는 base64url 인코딩 — padding 보정
            padded = data_b64.replace("-", "+").replace("_", "/")
            padded += "=" * (4 - len(padded) % 4)
            raw = base64.b64decode(padded)
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(raw)
                tmp_path = tmp.name
            text = _extract_file_text(tmp_path, ext)
            os.unlink(tmp_path)
            if text.strip():
                texts.append(f"[{fname}]\n{text[:3000]}")
        except Exception as e:
            texts.append(f"[{fname}] (추출 실패: {e})")
    return "\n\n".join(texts)[:5000]


def _extract_file_text(path, ext):
    try:
        if ext == ".pdf":
            r = subprocess.run(["pdftotext", "-q", path, "-"],
                               capture_output=True, text=True, timeout=60)
            return r.stdout
        elif ext == ".docx":
            import docx
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            parts = [shape.text for slide in prs.slides
                     for shape in slide.shapes
                     if hasattr(shape, "text") and shape.text.strip()]
            return "\n".join(parts)
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for sn in wb.sheetnames:
                ws = wb[sn]
                rows.append(f"[{sn}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None and str(c).strip()]
                    if cells:
                        rows.append(" | ".join(cells))
            wb.close()
            return "\n".join(rows)
    except Exception:
        pass
    return ""


# ---------------------------------------------------------------------------
# v5.1 — NAS Qdrant 검색
# ---------------------------------------------------------------------------

def search_nas_qdrant(query_text, top_k=5):
    """nomic-embed-text로 쿼리 임베딩 → nas_ra_docs 검색 → [{path,filename,score,excerpt}]"""
    try:
        result = _post_json(OLLAMA_EMBED_URL, {
            "model": "nomic-embed-text",
            "prompt": query_text
        }, timeout=60)
        vector = result["embedding"]

        hits = _post_json(NAS_QDRANT_SEARCH, {
            "vector": vector,
            "limit": top_k,
            "with_payload": True,
            "score_threshold": 0.5
        }, timeout=30).get("result", [])

        return [{
            "path": h["payload"]["file_path"],
            "filename": h["payload"]["filename"],
            "score": round(h["score"], 3),
            "excerpt": h["payload"].get("text", "")[:300]
        } for h in hits]
    except Exception:
        return []


# ---------------------------------------------------------------------------
# v5.1 — 품질 점수 계산
# ---------------------------------------------------------------------------

def calc_quality_score(analysis_dict):
    """0-10 점수 계산 (action 길이, 체크리스트 항목, NAS 참조, JSON 완성도)"""
    score = 0
    action = analysis_dict.get("action", "")
    if len(action) > 80:
        score += 2
    elif len(action) >= 30:
        score += 1

    wp = analysis_dict.get("wp_comment", "")
    items = len([l for l in wp.split("\n") if re.match(r"^\s*(\d+\.|[-*])\s", l)])
    if items >= 3:
        score += 2
    elif items >= 2:
        score += 1

    refs = len(analysis_dict.get("nas_refs", []))
    if refs >= 2:
        score += 2
    elif refs >= 1:
        score += 1

    required = ["summary", "org", "region", "task_type", "deadline", "action", "priority"]
    if all(analysis_dict.get(f) for f in required):
        score += 2

    return min(int(score * 1.25), 10)




# ========================================================================
# v5.2 추가 — Codex & Copilot 호출 (3-Model Support)
# ========================================================================

def call_codex(prompt, max_tokens=2000):
    """OpenAI GPT-4 via localhost:5055 gateway (Hermes OAuth Gateway)"""
    try:
        payload = json.dumps({
            "model": "gpt-4",
            "messages": [
                {
                    "role": "system",
                    "content": "당신은 의료기기 규제/인허가(RA) 전문가입니다. JSON만 출력하세요."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "temperature": 0.1,
            "max_tokens": max_tokens
        }).encode()

        req = urllib.request.Request(
            "http://localhost:5055/v1/chat/completions",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": "Bearer sk-hermes-n8n"
            },
            method="POST"
        )
        
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read().decode())
            if data and "choices" in data:
                return data["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"[warn] Codex error: {e}", flush=True)
    return None


def call_copilot(prompt, max_tokens=2000):
    """Copilot Pro (Claude Sonnet) via localhost:5055 gateway"""
    try:
        payload = json.dumps({
            "model": "claude-sonnet-4",
            "messages": [
                {
                    "role": "system",
                    "content": "당신은 의료기기 규제/인허가(RA) 전문가입니다. JSON만 출력하세요."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "temperature": 0.1,
            "max_tokens": max_tokens
        }).encode()

        req = urllib.request.Request(
            "http://localhost:5055/v1/chat/completions",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": "Bearer sk-hermes-n8n"
            },
            method="POST"
        )
        
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read().decode())
            if data and "choices" in data:
                return data["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"[warn] Copilot error: {e}", flush=True)
    return None


def analyze_mail_v5_2(from_addr, subject, body, attachments="", attachment_files=None):
    """
    v5.2 — 3-Model Support (Codex, Copilot, GLM)
    동일한 NAS RAG 검색 결과를 기반으로 3개 모델 분석
    """
    import threading

    error_base = {
        "summary": "", "org": "", "region": "", "task_type": "",
        "deadline": None, "action": "", "priority": "medium",
        "attachments_note": None, "wp_comment": "", "nas_refs": [],
        "models": {"codex": None, "copilot": None, "glm": None}
    }

    # 1. 첨부파일 텍스트 추출
    attachment_text = ""
    if attachment_files:
        attachment_text = extract_attachment_text(attachment_files)

    # 2. NAS RAG 검색 (기본 + Intended Use + Design Validation 특화 쿼리)
    query = f"{subject} {body[:500]} {attachment_text[:500]}"
    nas_refs = search_nas_qdrant(query, top_k=5)

    # Issue #2: Intended Use / Indication for Use 전용 쿼리 강화
    intended_refs = search_nas_qdrant(
        f"Indication for Use Intended Use User Manual 사용 목적 1.2 {subject}",
        top_k=3
    )
    # Issue #4: Design Validation 카테고리 전용 탐색
    dv_refs = search_nas_qdrant(
        f"Design Validation 설계 검증 design verification HnX-P1 test report {subject}",
        top_k=3
    )
    # 중복 없이 병합 (path 기준)
    seen_paths = {r['path'] for r in nas_refs}
    for r in intended_refs + dv_refs:
        if r['path'] not in seen_paths:
            nas_refs.append(r)
            seen_paths.add(r['path'])

    # 3. 분석 프롬프트 생성 — Issue #3: 실제 NAS 파일 경로 포함
    nas_section = ""
    if nas_refs:
        nas_lines = [
            f"- 경로: `{r['path']}` (유사도: {r['score']:.3f})\n  내용: {r['excerpt'][:150]}"
            for r in nas_refs[:5]
        ]
        nas_section = "\n[NAS 관련 문서 — 실제 파일 경로]\n" + "\n".join(nas_lines)

    attach_section = ""
    if attachments:
        attach_section += f"\n[첨부파일 목록]\n{attachments}"
    if attachment_text:
        attach_section += f"\n[첨부파일 내용]\n{attachment_text[:2000]}"

    prompt = (
        "당신은 의료기기 규제/인허가(RA) 전문가입니다. "
        "아래 공문을 분석하여 JSON만 출력하세요. 다른 텍스트, 설명, 마크다운 없이 순수 JSON만.\n\n"
        f"[발신자] {from_addr}\n"
        f"[제목] {subject}\n"
        f"[본문]\n{clean_body(body)}"
        f"{attach_section}{nas_section}\n\n"
        "출력 형식 (모든 필드 필수):\n"
        '{{'
        '"summary": "핵심 내용 3-5줄 요약",'
        '"org": "발신기관명",'
        '"region": "EU/FDA/KR/글로벌/기타",'
        '"task_type": "지원사업/규제변경/인증갱신/심사/공고/기타",'
        '"deadline": "마감기한 또는 null",'
        '"action": "RA담당자가 즉시 취해야 할 구체적 조치 (80자 이상)",'
        '"priority": "high/medium/low",'
        '"attachments_note": "첨부파일 주요사항 또는 null"'
        '}}'
    )

    # 4. 3개 모델 병렬 호출
    models = {"codex": None, "copilot": None, "glm": None}

    def call_model(name, func):
        try:
            if name == "glm":
                raw = call_glm("glm-4.5-air", prompt, max_tokens=2000)
            elif name == "codex":
                raw = call_codex(prompt, max_tokens=2000)
            elif name == "copilot":
                raw = call_copilot(prompt, max_tokens=2000)
            else:
                raw = func(prompt, max_tokens=2000)
            
            print(f"[debug] {name} raw response: {len(raw) if raw else 0} chars, type={type(raw)}", flush=True)
            
            if raw:
                parsed = parse_json_response(raw)
                print(f"[debug] {name} parsed: {type(parsed).__name__ if parsed else 'None'}", flush=True)
                if parsed:
                    models[name] = parsed
                    print(f"[debug] {name} result stored", flush=True)
                else:
                    print(f"[warn] {name} parse_json_response returned None", flush=True)
            else:
                print(f"[warn] {name} returned empty/None", flush=True)
        except Exception as e:
            print(f"[error] {name} model failed: {e}", flush=True)

    # 스레드 생성 및 병렬 실행
    threads = [
        threading.Thread(target=call_model, args=("codex", call_codex)),
        threading.Thread(target=call_model, args=("copilot", call_copilot)),
        threading.Thread(target=call_model, args=("glm", lambda p, m: call_glm("glm-4.5-air", p, m)))
    ]

    for t in threads:
        t.start()

    for t in threads:
        t.join(timeout=180)

    # 5. 결과 정리 (3개 모델 결과 + NAS refs 모두 포함)
    result = error_base.copy()
    result["nas_refs"] = nas_refs
    result["models"] = models

    # Primary model 선택 (Codex > Copilot > GLM 순서)
    if models["codex"]:
        result.update(models["codex"])
        result["primary_model"] = "codex"
    elif models["copilot"]:
        result.update(models["copilot"])
        result["primary_model"] = "copilot"
    elif models["glm"]:
        result.update(models["glm"])
        result["primary_model"] = "glm"
    else:
        result["summary"] = "3개 모델 분석 모두 실패"
        result["primary_model"] = "none"

    # 6. wp_comment 생성 (기존 로직)
    try:
        result["wp_comment"] = _run_wp_comment(result, nas_refs, attachment_text)
    except Exception:
        result["wp_comment"] = ""

    return result


# ---------------------------------------------------------------------------
# v5.1 — GLM API 호출
# ---------------------------------------------------------------------------

def call_glm(model, prompt, max_tokens=3000):
    """z.ai GLM API (OpenAI-compatible). 실패 시 None 반환."""
    if not GLM_API_KEY:
        print(f"[warn] GLM skipped: GLM_API_KEY not set", flush=True)
        return None
    payload = json.dumps({
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": max_tokens,
        "temperature": 0.1
    }).encode()
    req = urllib.request.Request(
        GLM_BASE_URL, data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {GLM_API_KEY}"
        }, method="POST"
    )
    try:
        print(f"[debug] GLM calling {model}...", flush=True)
        with urllib.request.urlopen(req, timeout=180) as resp:
            data = json.loads(resp.read())
            result = data["choices"][0]["message"]["content"]
            print(f"[debug] GLM success: {len(result)} chars", flush=True)
            return result
    except Exception as e:
        print(f"[error] GLM failed: {type(e).__name__}: {str(e)[:200]}", flush=True)
        return None


def _build_analysis_prompt(from_addr, subject, body, attachments, attachment_text):
    attach_section = ""
    if attachments:
        attach_section += f"\n[첨부파일 목록]\n{attachments}"
    if attachment_text:
        attach_section += f"\n[첨부파일 내용]\n{attachment_text[:2000]}"
    return ANALYSIS_PROMPT_TMPL.format(
        from_addr=from_addr,
        subject=subject,
        body=clean_body(body),
        attach_section=attach_section
    )


def _run_ollama_analyze(prompt):
    # GLM API 우선 사용 (RPi5 Ollama gemma3:4b 너무 느림)
    if GLM_API_KEY:
        result = call_glm("glm-4.5-air", prompt, max_tokens=2000)
        if result:
            return result
    result = _post_json(OLLAMA_GENERATE_URL, {
        "model": "gemma3:4b",
        "prompt": prompt,
        "stream": False,
        "options": {"temperature": 0.1, "num_predict": 800}
    }, timeout=120)
    return result.get("response", "")


def _run_wp_comment(analysis_dict, nas_refs, attachment_text):
    # NAS 문서: 상위 3개만, 각 발췌 200자로 제한 (프롬프트 길이 최적화)
    nas_text = "\n".join(
        f"- `{r['path']}` (score: {r['score']:.2f})\n  {r['excerpt'][:200]}"
        for r in (nas_refs[:3] if nas_refs else [])
    ) if nas_refs else "관련 NAS 문서 없음"

    # 프롬프트 최적화: JSON 전체 대신 핵심 정보만 압축 형식으로 (토큰 수 1/4 감소)
    analysis_compact = (
        f"업무: {analysis_dict.get('task_type', '미분류')}\n"
        f"기관: {analysis_dict.get('org', '')}\n"
        f"지역: {analysis_dict.get('region', '')}\n"
        f"기한: {analysis_dict.get('deadline', 'N/A')}\n"
        f"우선순위: {analysis_dict.get('priority', '')}\n"
        f"요약: {analysis_dict.get('summary', '')[:400]}"
    )
    att_text = (attachment_text or "없음")[:800]

    prompt = (
        "의료기기 RA 전담 에이전트로서 업무 체크리스트를 마크다운으로 작성하세요.\n\n"
        "[정보]\n" + analysis_compact + "\n\n"
        "[관련 문서]\n" + nas_text + "\n\n"
        "[첨부]\n" + att_text + "\n\n"
        "형식:\n"
        "## 🤖 Hermes RA 가이드\n\n"
        "### 분석\n(2줄)\n\n"
        "### 체크리스트\n1. 항목\n2. 항목\n3. 항목\n\n"
        "### 문서\n- 파일경로 — 발췌"
    )
    if GLM_API_KEY:
        raw = call_glm("glm-4.5-air", prompt, max_tokens=3000)
        if raw:
            return re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL).strip()
        print(f"[warn] GLM wp_comment 실패 (프롬프트: {len(prompt)} chars)", flush=True)
    return "## Hermes RA 가이드\n\n(wp_comment 생성 실패 — 분석 결과 참고)"




def analyze_mail_v5(from_addr, subject, body, attachments="", attachment_files=None):
    """
    v5.1 전체 파이프라인:
    1. 첨부파일 텍스트 추출
    2. Qdrant nas_ra_docs 검색
    3. gemma3:4b 1차 분석 → 품질 평가 → GLM 캐스케이드
    4. gemma3:4b 2차 wp_comment 생성
    """
    error_base = {
        "summary": "", "org": "", "region": "", "task_type": "",
        "deadline": None, "action": "", "priority": "medium",
        "attachments_note": None, "wp_comment": "", "nas_refs": []
    }

    # 1. 첨부파일 텍스트 추출
    attachment_text = ""
    if attachment_files:
        attachment_text = extract_attachment_text(attachment_files)

    # 2. NAS 시맨틱 검색
    query = f"{subject} {body[:500]} {attachment_text[:500]}"
    nas_refs = search_nas_qdrant(query, top_k=5)

    # 3. 1차 분석 (gemma3:4b)
    prompt = _build_analysis_prompt(from_addr, subject, body, attachments, attachment_text)
    try:
        raw1 = _run_ollama_analyze(prompt)
        result = parse_json_response(raw1)
    except Exception as e:
        error_base["summary"] = f"분석오류: {e}"
        return error_base

    if not result:
        error_base["summary"] = "JSON파싱실패"
        return error_base

    result["nas_refs"] = nas_refs

    # 4. 품질 기반 GLM 캐스케이드
    score = calc_quality_score(result)
    if score < 8 and GLM_API_KEY:
        if score < 5:
            glm_raw = call_glm("glm-5.1", prompt, max_tokens=3000)
        else:
            glm_raw = call_glm("glm-4.5-air", prompt, max_tokens=2000)
            if glm_raw:
                glm_res = parse_json_response(glm_raw)
                if glm_res:
                    glm_res["nas_refs"] = nas_refs
                    if calc_quality_score(glm_res) >= 7:
                        result = glm_res
                        glm_raw = None
                    else:
                        glm_raw = call_glm("glm-5.1", prompt, max_tokens=3000)
        if glm_raw:
            glm_res = parse_json_response(glm_raw)
            if glm_res:
                glm_res["nas_refs"] = nas_refs
                result = glm_res

    # 5. 2차 wp_comment 생성
    try:
        result["wp_comment"] = _run_wp_comment(result, nas_refs, attachment_text)
    except Exception:
        result["wp_comment"] = ""

    if "nas_refs" not in result:
        result["nas_refs"] = nas_refs

    return result


# ---------------------------------------------------------------------------
# v5 — /qa RAG (기존 유지)
# ---------------------------------------------------------------------------

DOMAIN_TERMS = [
    "EUDAMED", "DUNS", "PPWR", "MFDS", "UDI", "GMP", "MDR", "MDD", "PMCF", "PMS",
    "GUDID", "Javitech", "Annual Establishment", "Small Business", "Pre-sub", "PRE-SUB",
    "510k", "CE 인증", "CE 사후심사",
]


def extract_keywords(question):
    found = []
    q_lower = question.lower()
    for term in DOMAIN_TERMS:
        if term.lower() in q_lower:
            found.append(term)
    if "fda" in q_lower:
        if "annual" in q_lower or "establishment" in q_lower:
            found.append("Annual Establishment")
        elif "510k" in q_lower or "pre-sub" in q_lower:
            found.append("510k")
        else:
            found.append("FDA")
    return list(set(found))


def scroll_qdrant_keyword(keywords, limit=300):
    if not keywords:
        return []
    try:
        result = _post_json(QA_QDRANT_SCROLL_URL, {
            "limit": limit, "with_payload": True, "with_vectors": False
        }, timeout=30)
        points = result.get("result", {}).get("points", [])
        hits = [
            p for p in points
            if any(kw.upper() in p.get("payload", {}).get("text", "").upper() for kw in keywords)
        ]
        for h in hits:
            h["score"] = 0.75
            h["is_keyword_hit"] = True
        return hits
    except Exception:
        return []


DOC_TYPE_PRIORITY = {
    "handover_item": 0, "work_progress": 1, "annual_task": 2,
    "project_item": 3, "progress_status": 4, "certification_status": 5,
    "account_info": 9, "regulatory_knowledge": 9, "qa_history": 10,
}


def format_keyword_answer(hits, question):
    sorted_hits = sorted(hits, key=lambda h: DOC_TYPE_PRIORITY.get(
        h.get("payload", {}).get("doc_type", ""), 8))
    lines = []
    for h in sorted_hits[:4]:
        text = h.get("payload", {}).get("text", "").strip()
        doc_type = h.get("payload", {}).get("doc_type", "")
        if text and doc_type not in ("account_info", "qa_history"):
            lines.append(f"[{doc_type}] {text}")
    if not lines:
        for h in sorted_hits[:3]:
            text = h.get("payload", {}).get("text", "").strip()
            if text:
                lines.append(text)
    return "\n\n".join(lines) if lines else None


def embed_question(question):
    result = _post_json(OLLAMA_EMBED_URL, {
        "model": "nomic-embed-text",
        "prompt": question
    }, timeout=60)
    return result["embedding"]


STRUCTURED_DOC_TYPES = ["annual_task", "handover_item", "work_progress", "project_item"]


def search_qdrant(vector, top_k, score_threshold=0.0, filter_payload=None):
    payload = {"vector": vector, "limit": top_k,
               "with_payload": True, "score_threshold": score_threshold}
    if filter_payload:
        payload["filter"] = filter_payload
    result = _post_json(QA_QDRANT_SEARCH_URL, payload, timeout=30)
    return result.get("result", [])


def call_llm_rag(context, question):
    result = _post_json(OLLAMA_GENERATE_URL, {
        "model": "gemma3:4b",
        "system": (
            "당신은 의료기기 RA(규제 인허가) 전문 어시스턴트입니다. "
            "제공된 Context에서만 정보를 찾아 한국어로 답변하세요. "
            "날짜, 비용, 업체명, 담당자 등 구체적 정보를 최우선으로 추출하여 답하세요. "
            "Context에 관련 정보가 있으면 반드시 구체적으로 인용하세요. "
            "정말로 없을 때만 '관련 정보를 찾을 수 없습니다'라고 하세요."
        ),
        "prompt": f"Context:\n{context}\n\nQuestion: {question}\n\nAnswer:",
        "stream": False,
        "options": {"temperature": 0.1, "num_predict": 100, "num_ctx": 1024}
    }, timeout=600)
    return result.get("response", "").strip()


def handle_qa(question, top_k=5):
    keywords = extract_keywords(question)
    if keywords:
        kw_hits = scroll_qdrant_keyword(keywords)
        if kw_hits:
            answer = format_keyword_answer(kw_hits, question)
            if answer:
                sources = [{
                    "sheet": h.get("payload", {}).get("sheet", ""),
                    "doc_type": h.get("payload", {}).get("doc_type", ""),
                    "score": 0.75,
                    "excerpt": h.get("payload", {}).get("text", "")[:200]
                } for h in kw_hits[:6]]
                return {"answer": answer, "sources": sources, "question": question}, 200

    try:
        vector = embed_question(question)
    except Exception as e:
        return {"error": "service unavailable", "detail": str(e)}, 503

    structured_filter = {
        "should": [{"key": "doc_type", "match": {"value": dt}} for dt in STRUCTURED_DOC_TYPES]
    }
    try:
        structured_hits = search_qdrant(vector, top_k, filter_payload=structured_filter)
        general_hits = search_qdrant(vector, top_k)
    except Exception as e:
        return {"error": "service unavailable", "detail": str(e)}, 503

    seen_ids = set()
    hits = []
    for h in structured_hits + general_hits:
        hid = h.get("id")
        if hid not in seen_ids:
            seen_ids.add(hid)
            hits.append(h)
    hits = hits[:min(top_k, 5)]

    if not hits:
        return {"answer": "No relevant information found in knowledge base.", "sources": [], "question": question}, 200

    context_parts = []
    sources = []
    for i, hit in enumerate(hits):
        payload = hit.get("payload", {})
        text = payload.get("text", "")
        doc_type = payload.get("doc_type", "")
        sheet = payload.get("sheet", "")
        context_parts.append(f"[출처 {i+1}] ({doc_type}/{sheet})\n{text}")
        sources.append({"sheet": sheet, "doc_type": doc_type,
                        "score": hit.get("score", 0.0), "excerpt": text[:200]})
    context = "\n\n".join(context_parts)

    try:
        answer = call_llm_rag(context, question)
    except Exception as e:
        return {"error": "service unavailable", "detail": str(e)}, 503

    return {"answer": answer, "sources": sources, "question": question}, 200


def handle_health():
    ollama_ok = _get_json(OLLAMA_TAGS_URL, timeout=5) is not None
    qdrant_data = _get_json(QA_QDRANT_COLLECTION_URL, timeout=5)
    qdrant_ok = qdrant_data is not None
    kb_points = 0
    if qdrant_ok and isinstance(qdrant_data, dict):
        result_obj = qdrant_data.get("result", {})
        kb_points = result_obj.get("points_count") or result_obj.get("vectors_count") or 0
    status = "ok" if (ollama_ok and qdrant_ok) else "degraded"
    return {"status": status, "ollama": ollama_ok, "qdrant": qdrant_ok, "kb_points": kb_points}, 200


# ---------------------------------------------------------------------------
# HTTP 핸들러
# ---------------------------------------------------------------------------

class Handler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        pass

    def _send_json(self, status, payload):
        resp = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", len(resp))
        self.end_headers()
        self.wfile.write(resp)

    def do_GET(self):
        if self.path == "/health":
            result, status = handle_health()
            self._send_json(status, result)
        elif self.path == "/check-github-token":
            import subprocess as _sp
            r = _sp.run(["gh", "auth", "status"], capture_output=True, timeout=10)
            ok = r.returncode == 0
            self._send_json(200, {"token_ok": ok,
                                   "status": (r.stdout + r.stderr).decode(errors="replace").strip()[:200]})
        else:
            self._send_json(404, {"error": "not found"})

    def do_POST(self):
        length = int(self.headers.get("Content-Length", 0))

        if self.path == "/analyze":
            try:
                data = json.loads(self.rfile.read(length).decode("utf-8"))
                result = analyze_mail_v5_2(
                    data.get("from", ""),
                    data.get("subject", ""),
                    data.get("body", ""),
                    data.get("attachments", ""),
                    data.get("attachment_files", None)
                )
                self._send_json(200, result)
            except Exception as e:
                self._send_json(500, {"error": str(e)})

        elif self.path == "/qa":
            try:
                data = json.loads(self.rfile.read(length).decode("utf-8"))
            except Exception as e:
                self._send_json(400, {"error": "invalid JSON", "detail": str(e)})
                return
            question = data.get("question", "").strip()
            if not question:
                self._send_json(400, {"error": "question field is required"})
                return
            top_k = data.get("top_k", 5)
            try:
                top_k = int(top_k)
            except (TypeError, ValueError):
                top_k = 5
            result, status = handle_qa(question, top_k)
            self._send_json(status, result)

        elif self.path == "/run-extract-mail":
            import threading
            try:
                data = json.loads(self.rfile.read(length).decode("utf-8"))
                trigger = data.get("trigger", "manual")
                file_path = data.get("file", "")
            except Exception:
                trigger, file_path = "manual", ""

            def run_extract():
                import subprocess as _sp2
                r2 = _sp2.run(["python3", "/opt/hermes/extract_mail_qa.py"],
                              capture_output=True, timeout=3600)
                with open("/tmp/extract_mail_qa.log", "a") as f:
                    f.write(f"[auto-trigger] trigger={trigger} file={file_path}\n")
                    if r2.stdout:
                        f.write(r2.stdout.decode(errors="replace"))
                    if r2.stderr:
                        f.write(r2.stderr.decode(errors="replace"))

            import subprocess as _sp
            running = _sp.run(["pgrep", "-f", "extract_mail_qa.py"], capture_output=True).returncode == 0
            if running:
                self._send_json(200, {"status": "already_running", "trigger": trigger})
            else:
                t = threading.Thread(target=run_extract, daemon=True)
                t.start()
                self._send_json(200, {"status": "started", "trigger": trigger, "file": file_path})

        else:
            self._send_json(404, {"error": "not found"})


class ThreadingHTTPServer(ThreadingMixIn, HTTPServer):
    daemon_threads = True


if __name__ == "__main__":
    server = ThreadingHTTPServer(("0.0.0.0", 7788), Handler)
    print("Hermes RA API server v5.2 (3-Model Support) running on port 7788")
    server.serve_forever()
